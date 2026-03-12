"""Artifacts module Pydantic models"""

from typing import Annotated, List, Literal, Union

from pptx.util import Inches, Pt
from pydantic import BaseModel, Field


class DownloadArtifactRequest(BaseModel):
    """成果物ダウンロードのリクエストモデル"""

    format: str = Field(
        ..., description="成果物のファイル形式(拡張子)", examples=["pptx"]
    )
    artifact_id: str = Field(
        ..., description="成果物のID", examples=["3a92c5d1-6f3b-4b9a-bc2c-9f7d6e4f8e7a"]
    )
    version: int = Field(..., description="成果物のバージョン", examples=[1])


class EncodedArtifact(BaseModel):
    """Base64エンコードされた成果物のレスポンスモデル"""

    format: str = Field(..., description="ファイル形式", examples=["pptx"])
    content: str = Field(..., description="Base64エンコードされたファイル内容")


# -----LLMによるパースによってマークダウンから変換されるPowerPointオブジェクト-----


# ページの基底クラス
class Page(BaseModel):
    header: str = Field(..., description="スライドのヘッダー情報。ページのタイトル。")
    content: str = Field(..., description="スライドに含むコンテンツ概要")
    template: Literal[
        "text",
        "table",
    ] = Field(
        ...,
        description="スライドのコンテンツタイプ\n"
        "text: テキストのみのページ。複数セクションと複数個の箇条書きで構成\n"
        "table: 表のみのページ。表のみで構成\n",
    )
    policy: str = Field(
        ..., description="スライド作成にあたるデザインやコンテンツの作成方針を記載"
    )


# テキストセクションの定義
class TextSection(BaseModel):
    title: str = Field(..., description="テキストセクションのタイトル")
    content: Annotated[
        list[str],
        Field(..., description="テキストセクションの内容。箇条書き形式"),
    ]


# テキストページ
class TextPage(Page):
    sections: Annotated[
        list[TextSection],
        Field(..., description="スライドのコンテンツセクションのリスト"),
    ]


# 表セクションの定義
class TableSection(BaseModel):
    title: str = Field(..., description="表セクションのタイトル（オプション）")
    table_data: List[List[str]] = Field(..., description="表データ（二次元リスト）")


# 表ページ
class TablePage(Page):
    table_section: Annotated[
        TableSection,
        Field(..., description="スライドの表セクション"),
    ]


# パワーポイント全体
class PowerPoint(BaseModel):
    title: str = Field(..., description="プレゼンテーションのタイトル")
    pages: List[
        Union[
            TextPage,
            TablePage,
        ]
    ] = Field(..., description="プレゼンテーションの各ページ情報のリスト")


# -----PowerPointオブジェクトから設定されるPowerPointプレゼンテーションのスライドオブジェクト-----


# スライドの基底クラス
class SlideBase:
    def __init__(self, presentation, page_data):
        self.presentation = presentation
        self.page_data = page_data
        self.slide = None

    def set_title(self, title_text):
        """タイトルを設定する"""
        title = self.slide.placeholders[0]
        title.text = title_text


# タイトルスライドのクラス
class TitleSlide(SlideBase):
    def create_slide(self):
        """タイトルスライドを作成する"""
        slide_layout = self.presentation.slide_layouts[0]
        self.slide = self.presentation.slides.add_slide(slide_layout)
        self.set_title(self.page_data)


# テキストスライドのクラス
class TextSlide(SlideBase):
    def create_slide(self):
        """テキストスライドを作成する"""
        slide_layout = self.presentation.slide_layouts[1]
        self.slide = self.presentation.slides.add_slide(slide_layout)
        self.set_title(self.page_data.header)

        content = self.slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()

        self.add_sections_to_slide(text_frame, self.page_data.sections)

    def add_sections_to_slide(self, text_frame, sections):
        """箇条書きを作成する"""
        for section in sections:
            p = text_frame.add_paragraph()
            p.text = section.title
            p.level = 0

            for item in section.content:
                p = text_frame.add_paragraph()
                p.text = f"- {item}"
                p.level = 1


# 表スライドのクラス
class TableSlide(SlideBase):
    def create_slide(self):
        """表スライドを作成する"""
        slide_layout = self.presentation.slide_layouts[1]
        self.slide = self.presentation.slides.add_slide(slide_layout)
        self.set_title(self.page_data.header)

        # テキストボックスを削除
        placeholder = self.slide.placeholders[1]
        sp = placeholder._sp
        sp.getparent().remove(sp)

        # 表の設定
        left = Inches(0.78740157)
        top = Inches(1.37795276)
        width = Inches(11.7519685)
        height = Inches(5.33464567)

        font_size = Pt(16)

        # 表の挿入
        table = self.slide.shapes.add_table(
            rows=len(self.page_data.table_section.table_data),
            cols=len(self.page_data.table_section.table_data[0]),
            left=left,
            top=top,
            width=width,
            height=height,
        ).table

        # 表のデータを設定
        for row_index, row_data in enumerate(self.page_data.table_section.table_data):
            for col_index, cell_data in enumerate(row_data):
                cell = table.cell(row_index, col_index)
                cell.text = str(cell_data)

                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = font_size
